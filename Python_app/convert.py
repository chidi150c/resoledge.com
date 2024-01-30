from pptx import Presentation
from PIL import Image

# Load the PowerPoint file
ppt = Presentation('your_presentation.pptx')

# Create a directory to save the images (if it doesn't exist)
import os
if not os.path.exists('ppt_images'):
    os.makedirs('ppt_images')

# Extract and save each slide as an image
for i, slide in enumerate(ppt.slides):
    image = slide.shapes.add_picture('ppt_images/slide{}.png'.format(i), 0, 0, width=960, height=540)
    image = Image.open('ppt_images/slide{}.png'.format(i))
    image.save('ppt_images/slide{}.png'.format(i), 'PNG')

print('Slides extracted and saved as images.')
